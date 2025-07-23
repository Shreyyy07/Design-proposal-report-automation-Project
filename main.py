import streamlit as st
import pandas as pd
import fitz  # PyMuPDF for enhanced PDF processing
import io
from PIL import Image as PILImage
import tempfile
import os
import subprocess
import pyautogui
import pygetwindow as gw
import time
import shutil
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import textwrap
import glob
import winreg
import numpy as np

# ---------- ASSET PATHS ----------
# Define paths for assets to make them easy to change and manage.
ASSETS_DIR = "assets"
LOGO_PATH = os.path.join(ASSETS_DIR, "apollo_logo.png")
SLIDE1_PATH = os.path.join(ASSETS_DIR, "slide1.png")
PD_PATH = os.path.join(ASSETS_DIR, "pd.png")
PD2_PATH = os.path.join(ASSETS_DIR, "pd2.png")
MOM_PATH = os.path.join(ASSETS_DIR, "mom.png")
LASTSLIDE_PATH = os.path.join(ASSETS_DIR, "lastslide.png")
HEADER_BANNER_PATH = os.path.join(ASSETS_DIR, "header_banner.png")

def get_apollo_logo_base64():
    """
    Reads the Apollo Tyres logo image and converts it to a base64 string.
    """
    try:
        with open(LOGO_PATH, "rb") as img_file:
            return base64.b64encode(img_file.read()).decode()
    except FileNotFoundError:
        st.warning(f"Logo not found at {LOGO_PATH}. Please ensure it exists in the 'assets' folder.")
        return None

def add_logo_to_streamlit():
    """
    Displays the Apollo Tyres logo at the top-center of the Streamlit application.
    """
    logo_base64 = get_apollo_logo_base64()
    if logo_base64:
        st.markdown(
            f"""
            <div style="display: flex; justify-content: center; margin-bottom: 20px;">
                <img src="data:image/png;base64,{logo_base64}" width="200" alt="Apollo Tyres Logo">
            </div>
            """,
            unsafe_allow_html=True
        )
    else:
        st.markdown(
            """
            <div style="text-align: center; margin-bottom: 20px;">
                <h2 style="color: #8A2BE2; margin: 0;">🏎️ APOLLO TYRES LTD</h2>
            </div>
            """,
            unsafe_allow_html=True
        )

# ---------- SOFTWARE DETECTION FUNCTIONS ----------

def find_autocad_executable():
    """
    Automatically finds the AutoCAD executable ('acad.exe') on the system.
    """
    if 'custom_autocad_path' in st.session_state and st.session_state['custom_autocad_path']:
        custom_path = st.session_state['custom_autocad_path']
        if os.path.exists(custom_path):
            return custom_path
    
    possible_paths = [
        "C:/Program Files/Autodesk/AutoCAD 2025/acad.exe",
        "C:/Program Files/Autodesk/AutoCAD 2024/acad.exe",
        "C:/Program Files/Autodesk/AutoCAD 2023/acad.exe",
    ]
    for path in possible_paths:
        if os.path.exists(path):
            return path
    
    autocad_patterns = ["C:/Program Files/Autodesk/AutoCAD*/acad.exe"]
    for pattern in autocad_patterns:
        matches = glob.glob(pattern)
        if matches:
            matches.sort(reverse=True)
            return matches[0]
    
    return None

def find_nx_executable():
    """
    Automatically finds the Siemens NX executable ('ugraf.exe') on the system.
    """
    if 'custom_nx_path' in st.session_state and st.session_state['custom_nx_path']:
        custom_path = st.session_state['custom_nx_path']
        if os.path.exists(custom_path):
            return custom_path
    
    nx_patterns = [
        "C:/Program Files/Siemens/NX*/NXBIN/ugraf.exe",
        "D:/abcde/NXBIN/ugraf.exe",
    ]
    
    for pattern in nx_patterns:
        matches = glob.glob(pattern)
        if matches:
            matches.sort(reverse=True)
            return matches[0]
            
    return None

# ---------- SESSION STATE INITIALIZATION ----------

def initialize_session_state():
    """
    Initializes the Streamlit session state variables.
    """
    if 'cad_screenshots_captured' not in st.session_state:
        st.session_state['cad_screenshots_captured'] = False
    if 'cad_screenshot_paths' not in st.session_state:
        st.session_state['cad_screenshot_paths'] = []
    if 'nx_screenshots_captured' not in st.session_state:
        st.session_state['nx_screenshots_captured'] = False
    if 'nx_model_groups' not in st.session_state:
        st.session_state['nx_model_groups'] = []
    if 'custom_autocad_path' not in st.session_state:
        st.session_state['custom_autocad_path'] = ""
    if 'custom_nx_path' not in st.session_state:
        st.session_state['custom_nx_path'] = ""

# ---------- PDF EXTRACTION ----------
def extract_pdf_elements(pdf_path):
    """
    Extracts both text and images from a PDF, page by page, with their bounding boxes.
    """
    doc = fitz.open(pdf_path)
    pages_elements = []
    
    for page_num, page in enumerate(doc):
        page_elements = []
        blocks = page.get_text("dict", flags=11)["blocks"]
        for b in blocks:
            for l in b["lines"]:
                for s in l["spans"]:
                    page_elements.append({
                        "type": "text",
                        "bbox": s["bbox"],
                        "text": s["text"],
                        "size": s["size"] # Also extract font size
                    })

        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            
            img_rects = page.get_image_rects(img)
            if img_rects:
                bbox = img_rects[0]
                page_elements.append({
                    "type": "image",
                    "bbox": bbox,
                    "bytes": image_bytes
                })

        page_elements.sort(key=lambda el: (el["bbox"][1], el["bbox"][0]))
        pages_elements.append((page.rect, page_elements))
        
    return pages_elements


def extract_pdf_elements_improved(pdf_path):
    """
    Improved PDF extraction that better handles text structure and images.
    """
    doc = fitz.open(pdf_path)
    pages_elements = []
    
    for page_num, page in enumerate(doc):
        page_elements = {"text_blocks": [], "images": []}
        
        # Extract text blocks (preserves paragraph structure)
        text_dict = page.get_text("dict")
        for block in text_dict["blocks"]:
            if "lines" in block:
                block_text = ""
                for line in block["lines"]:
                    line_text = ""
                    for span in line["spans"]:
                        line_text += span["text"]
                    if line_text.strip():
                        block_text += line_text.strip() + " "
                
                if block_text.strip():
                    page_elements["text_blocks"].append({
                        "text": block_text.strip(),
                        "bbox": block["bbox"],
                        "block_type": "paragraph"
                    })
        
        # Extract images with better positioning
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Get image rectangles for positioning
                img_rects = page.get_image_rects(img)
                if img_rects:
                    bbox = img_rects[0]
                    page_elements["images"].append({
                        "bytes": image_bytes,
                        "bbox": bbox,
                        "index": img_index
                    })
            except Exception as e:
                print(f"Error extracting image {img_index}: {e}")
        
        pages_elements.append((page.rect, page_elements))
    
    doc.close()
    return pages_elements

# ---------- EXCEL READING ----------
def read_excel_data(excel_path):
    """
    Reads data from an Excel file into a pandas DataFrame.
    """
    try:
        return pd.read_excel(excel_path)
    except Exception as e:
        st.error(f"Excel error: {e}")
        return pd.DataFrame()

# ---------- AUTOCAD FUNCTIONS (SMART CONTENT DETECTION) ----------
def open_autocad_and_capture_screenshot(dwg_path):
    """
    Automates opening AutoCAD and takes a precise screenshot with smart content detection
    to crop only the drawing area, removing all white space from sides.
    """
    autocad_process = None
    try:
        autocad_path = find_autocad_executable()
        if not autocad_path:
            st.error("❌ AutoCAD not found. Please configure the path in settings.")
            return None
        
        st.info(f"📍 Found AutoCAD at: {autocad_path}")
        autocad_process = subprocess.Popen([autocad_path, dwg_path])

        st.info("Opening AutoCAD... Please wait.")
        time.sleep(18)
        
        autocad_window = None
        for _ in range(10): 
            windows = gw.getWindowsWithTitle('AutoCAD')
            if windows:
                autocad_window = windows[0]
                break
            time.sleep(1)
        
        if not autocad_window:
            st.error("Could not find AutoCAD window. Automation failed.")
            if autocad_process: autocad_process.kill()
            return None

        autocad_window.activate()
        if not autocad_window.isMaximized:
            autocad_window.maximize()
        time.sleep(3)

        pyautogui.press('esc', presses=2, interval=0.3)
        pyautogui.typewrite('zoom\n')
        time.sleep(0.5)
        pyautogui.typewrite('e\n')
        time.sleep(5)  # Wait for zoom to complete

        # Additional wait time before taking screenshot
        st.info("Waiting for drawing to render properly...")
        time.sleep(3)

        win_left, win_top, win_width, win_height = autocad_window.left, autocad_window.top, autocad_window.width, autocad_window.height
        
        # Initial aggressive cropping to remove UI elements
        top_offset = 220
        bottom_offset = 140
        left_offset = 350
        right_offset = 120
        
        # Take initial screenshot of the viewport area
        initial_region = (
            win_left + left_offset,
            win_top + top_offset,
            win_width - left_offset - right_offset,
            win_height - top_offset - bottom_offset
        )
        
        initial_screenshot = pyautogui.screenshot(region=initial_region)
        
        # Convert to PIL Image for content detection
        img = initial_screenshot
        img_array = img.load()
        img_width, img_height = img.size
        
        # Find content boundaries by detecting non-white pixels
        # Threshold for "white" pixels (allowing slight variations)
        white_threshold = 240
        
        # Find leftmost non-white column
        left_content = 0
        for x in range(img_width):
            has_content = False
            for y in range(img_height):
                r, g, b = img_array[x, y][:3]
                if r < white_threshold or g < white_threshold or b < white_threshold:
                    has_content = True
                    break
            if has_content:
                left_content = max(0, x - 20)  # Add small margin
                break
        
        # Find rightmost non-white column
        right_content = img_width
        for x in range(img_width - 1, -1, -1):
            has_content = False
            for y in range(img_height):
                r, g, b = img_array[x, y][:3]
                if r < white_threshold or g < white_threshold or b < white_threshold:
                    has_content = True
                    break
            if has_content:
                right_content = min(img_width, x + 20)  # Add small margin
                break
        
        # Find topmost non-white row
        top_content = 0
        for y in range(img_height):
            has_content = False
            for x in range(img_width):
                r, g, b = img_array[x, y][:3]
                if r < white_threshold or g < white_threshold or b < white_threshold:
                    has_content = True
                    break
            if has_content:
                top_content = max(0, y - 10)  # Add small margin
                break
        
        # Find bottommost non-white row
        bottom_content = img_height
        for y in range(img_height - 1, -1, -1):
            has_content = False
            for x in range(img_width):
                r, g, b = img_array[x, y][:3]
                if r < white_threshold or g < white_threshold or b < white_threshold:
                    has_content = True
                    break
            if has_content:
                bottom_content = min(img_height, y + 10)  # Add small margin
                break
        
        # Crop to content boundaries
        if right_content > left_content and bottom_content > top_content:
            cropped_screenshot = img.crop((left_content, top_content, right_content, bottom_content))
        else:
            # Fallback: use original image if content detection fails
            cropped_screenshot = img
        
        output_image_path = os.path.join(tempfile.gettempdir(), "cad_screenshot.png")
        cropped_screenshot.save(output_image_path)

        return output_image_path

    except Exception as e:
        st.error(f"Error with AutoCAD: {e}")
        return None
    finally:
        if autocad_process:
            try:
                st.info("Terminating AutoCAD process...")
                autocad_process.terminate()
                autocad_process.wait(timeout=5)
            except Exception:
                autocad_process.kill()
    
# ---------- NX FUNCTIONS (INCREASED OPENING TIME) ----------
def open_nx_and_capture_views_manual():
    """
    Guides the user through capturing the three required 3D views from Siemens NX.
    Fixed to use unique filenames for each model capture.
    """
    try:
        nx_path = find_nx_executable()
        if not nx_path:
            st.error("❌ Siemens NX not found. Please configure the path in settings.")
            return {}
        st.info(f"📍 Found NX at: {nx_path}")
        subprocess.Popen([nx_path], shell=True)
        
        st.info("Opening NX... Please wait (this may take up to 30 seconds).")
        time.sleep(30)

        nx_window = gw.getWindowsWithTitle('NX')[0]
        nx_window.activate()
        if not nx_window.isMaximized:
            nx_window.maximize()
        time.sleep(2)

        st.warning("Please manually open your 3D file in NX now.")
        st.info("After opening, follow the instructions for each view.")
        time.sleep(10)
        
        views_to_capture = ['Top View', 'Front View', 'Isometric View']
        screenshots = {}
        status_placeholder = st.empty()
        progress_placeholder = st.empty()
        timer_placeholder = st.empty()

        # Get unique model number for this capture session
        model_count = len(st.session_state.get('nx_model_groups', []))
        
        for i, view in enumerate(views_to_capture):
            status_placeholder.info(f"📸 **Step {i+1}/{len(views_to_capture)}**: Please set the **{view}** in NX")
            progress_placeholder.progress((i) / len(views_to_capture))
            
            for t in range(5, 0, -1):
                timer_placeholder.info(f"Adjusting to {view}. Screenshot in {t} seconds...")
                time.sleep(1)
            
            timer_placeholder.empty()

            nx_window.activate()
            pyautogui.press('f')
            time.sleep(2)

            screenshot = pyautogui.screenshot()
            screen_width, screen_height = pyautogui.size()
            
            left = int(screen_width * 0.28)    
            top = int(screen_height * 0.20)
            right = int(screen_width * 0.75)   
            bottom = int(screen_height * 0.90)

            cropped_screenshot = screenshot.crop((left, top, right, bottom))
            centered_screenshot = center_tire_properly(cropped_screenshot)
            
            # Create unique filename with model number and timestamp
            timestamp = int(time.time())
            img_path = os.path.join(tempfile.gettempdir(), 
                                  f"nxview_model{model_count + 1}_{view.replace(' ', '_').lower()}_{timestamp}.png")
            centered_screenshot.save(img_path)
            screenshots[view] = img_path

            st.success(f"{view} captured!")

        status_placeholder.success("🎉 **ALL VIEWS CAPTURED!** 🎉")
        progress_placeholder.progress(1.0)
        st.info("Closing NX...")
        nx_window.close()
        time.sleep(2)
        try:
            pyautogui.press('n')
        except:
            pass
        
        return screenshots

    except Exception as e:
        st.error(f"Error in NX automation: {e}")
        return {}

# ---------- BETTER TIRE CENTERING ----------
def center_tire_properly(image):
    """
    Centers the tire properly by detecting tire and moving it to exact center.
    """
    try:
        img_array = np.array(image)
        
        if len(img_array.shape) == 3:
            gray = np.mean(img_array, axis=2)
        else:
            gray = img_array
        
        # Find tire pixels
        tire_mask = gray < 180
        
        if not np.any(tire_mask):
            return image
        
        y_coords, x_coords = np.where(tire_mask)
        
        if len(x_coords) == 0:
            return image
        
        # Find tire center
        tire_center_x = int(np.mean(x_coords))
        tire_center_y = int(np.mean(y_coords))
        
        img_height, img_width = gray.shape
        image_center_x = img_width // 2
        image_center_y = img_height // 2
        
        # Calculate movement needed
        shift_x = image_center_x - tire_center_x
        shift_y = image_center_y - tire_center_y
        
        # Create new centered image
        new_img = PILImage.new('RGB', (img_width, img_height), (200, 200, 200))
        
        # Calculate paste position
        paste_x = max(0, shift_x)
        paste_y = max(0, shift_y)
        
        # Calculate crop region from original
        crop_left = max(0, -shift_x)
        crop_top = max(0, -shift_y)
        crop_right = min(img_width, img_width - shift_x)
        crop_bottom = min(img_height, img_height - shift_y)
        
        if crop_right > crop_left and crop_bottom > crop_top:
            cropped = image.crop((crop_left, crop_top, crop_right, crop_bottom))
            new_img.paste(cropped, (paste_x, paste_y))
            return new_img
        
        return image
        
    except Exception:
        return image

# ---------- POWERPOINT GENERATION FUNCTIONS (UPDATED) ----------

def add_slide_banner(slide, title_text):
    """
    Adds a custom purple banner with a title to the top of a slide.
    """
    try:
        if os.path.exists(HEADER_BANNER_PATH):
             slide.shapes.add_picture(HEADER_BANNER_PATH, 0, 0, width=prs.slide_width, height=Inches(0.8))
        else:
            banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
            banner.fill.solid()
            banner.fill.fore_color.rgb = RGBColor(138, 43, 226)
            banner.line.fill.background()
    except:
        pass

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.33), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    title_frame.word_wrap = True

def create_outline_slide(prs, topics):
    """
    Creates the 'Outline' slide with a DYNAMIC bulleted list of the report sections.
    UPDATED: Now adds Apollo logo at bottom left corner.
    """
    blank_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(blank_layout)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1.0))
    title_tf = title_shape.text_frame
    p_title = title_tf.paragraphs[0]
    p_title.text = "Outline:-"
    p_title.font.size = Pt(44)
    p_title.font.bold = True

    body_shape = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11), Inches(5.0))
    tf = body_shape.text_frame
    tf.clear()

    for topic in topics:
        p = tf.add_paragraph()
        p.text = topic
        p.font.size = Pt(24)
        p.level = 0
        p.space_after = Pt(12)

    # Add Apollo logo at bottom left corner
    try:
        if os.path.exists(LOGO_PATH):
            logo_width = Inches(1.5)
            logo_height = Inches(0.75)
            logo_left = Inches(0.3)
            logo_top = prs.slide_height - logo_height - Inches(0.3)
            slide.shapes.add_picture(LOGO_PATH, logo_left, logo_top, logo_width, logo_height)
    except Exception as e:
        pass  # Logo addition failed, continue without it


def generate_ppt(pdf_pages, excel_df, cad_image_paths, nx_model_groups, output_path):
    """
    Main function to generate the entire PowerPoint presentation with the corrected slide order
    and dynamic outline.
    """
    global prs
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- 1. Front Page Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    try:
        slide1.shapes.add_picture(SLIDE1_PATH, 0, 0, prs.slide_width, prs.slide_height)
    except Exception:
        st.warning(f"Could not find {SLIDE1_PATH}. Slide 1 will be blank.")

    # --- 2. Dynamic Outline Slide ---
    outline_topics = []
    if pdf_pages:
        outline_topics.append("Project brief details")
    outline_topics.append("Minutes of the Meeting")
    if not excel_df.empty:
        outline_topics.append("Design input sheet")
    if cad_image_paths:
        outline_topics.append("Cavity and thread profile details")
    if nx_model_groups:
        outline_topics.append("3D model images")
        outline_topics.append("Front & isometric close-up view of tyre images")
    create_outline_slide(prs, outline_topics)

    # --- 3. pd.png Slide (RESTORED) ---
    slide2 = prs.slides.add_slide(blank_layout)
    try:
        slide2.shapes.add_picture(PD_PATH, 0, 0, prs.slide_width, prs.slide_height)
    except Exception:
        st.warning(f"Could not find {PD_PATH}. Slide 3 will be blank.")

    # --- 4. pd2.png Slide ---
    slide3 = prs.slides.add_slide(blank_layout)
    try:
        slide3.shapes.add_picture(PD2_PATH, 0, 0, prs.slide_width, prs.slide_height)
    except Exception:
        st.warning(f"Could not find {PD2_PATH}. Slide 4 will be blank.")
        
    # --- 5. Minutes of the Meeting Slide ---
    mom_slide = prs.slides.add_slide(blank_layout)
    add_slide_banner(mom_slide, "Minutes of the Meeting")
    try:
        img_width = Inches(12) 
        img_height = Inches(6)
        left = (prs.slide_width - img_width) / 2
        top = (prs.slide_height - img_height) / 2 + Inches(0.2)
        mom_slide.shapes.add_picture(MOM_PATH, left, top, width=img_width, height=img_height)
    except Exception:
        st.warning(f"Could not find {MOM_PATH}. MOM slide will be blank.")

    # --- Content Sections ---
    if pdf_pages:
        create_pdf_content_slides(prs, pdf_pages, section_title="Project brief details")
    if not excel_df.empty:
        create_excel_slides(prs, excel_df, section_title="Design input sheet:")
    if cad_image_paths:
        create_cad_slides(prs, cad_image_paths, section_title="Cavity and thread profile details")
    if nx_model_groups:
        create_nx_model_slides(prs, nx_model_groups, section_title="3D model images")
        # Add close-up view slide right after 3D model slides
        create_nx_closeup_slide(prs, nx_model_groups, section_title="Front & isometric close-up view of tyre images")

    # --- Final Thank You Slide ---
    thank_slide = prs.slides.add_slide(blank_layout)
    try:
        thank_slide.shapes.add_picture(LASTSLIDE_PATH, 0, 0, prs.slide_width, prs.slide_height)
    except Exception:
        st.warning(f"Could not find {LASTSLIDE_PATH}. Final slide will be blank.")

    prs.save(output_path)

def create_pdf_content_slides_improved(prs, pdf_pages, section_title):
    """
    Creates professional PDF slides with bullet points and proper image handling.
    """
    blank_layout = prs.slide_layouts[6]
    
    for page_idx, (page_rect, page_elements) in enumerate(pdf_pages):
        text_blocks = page_elements.get("text_blocks", [])
        images = page_elements.get("images", [])
        
        if not text_blocks and not images:
            continue
        
        # Process text into bullet points
        bullet_points = []
        for block in text_blocks:
            text = block["text"]
            
            # Split by common sentence endings and bullet indicators
            sentences = []
            for delimiter in ['. ', '• ', '◦ ', '- ', '\n', '• ', '·']:
                if delimiter in text:
                    parts = text.split(delimiter)
                    sentences.extend([part.strip() for part in parts if part.strip()])
                    break
            else:
                # If no delimiters found, split by length
                words = text.split()
                chunk_size = 15
                sentences = [' '.join(words[i:i + chunk_size]) 
                           for i in range(0, len(words), chunk_size)]
            
            bullet_points.extend([s for s in sentences if len(s) > 10])
        
        # Create slides with bullet points
        points_per_slide = 8
        num_text_slides = (len(bullet_points) + points_per_slide - 1) // points_per_slide if bullet_points else 0
        
        for slide_idx in range(max(1, num_text_slides)):
            slide = prs.slides.add_slide(blank_layout)
            
            # Add slide title
            if num_text_slides > 1:
                title = f"{section_title} - Page {page_idx + 1} (Part {slide_idx + 1})"
            else:
                title = f"{section_title} - Page {page_idx + 1}"
            
            add_slide_banner(slide, title)
            
            # Add bullet points
            if bullet_points:
                start_idx = slide_idx * points_per_slide
                end_idx = min(start_idx + points_per_slide, len(bullet_points))
                slide_points = bullet_points[start_idx:end_idx]
                
                # Create bullet point text box
                text_box = slide.shapes.add_textbox(
                    Inches(0.5),
                    Inches(1.2),
                    Inches(12.33),
                    Inches(5.8)
                )
                
                text_frame = text_box.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.NONE
                
                for i, point in enumerate(slide_points):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = point
                    p.font.size = Pt(16)
                    p.space_after = Pt(12)
                    p.line_spacing = 1.3
                    p.level = 0  # This creates bullet points
                    p.alignment = PP_ALIGN.LEFT
        
        # Create separate slides for images
        for img_idx, image_data in enumerate(images):
            slide = prs.slides.add_slide(blank_layout)
            add_slide_banner(slide, f"{section_title} - Page {page_idx + 1} - Image {img_idx + 1}")
            
            try:
                # Calculate image size and position
                page_width, page_height = page_rect.width, page_rect.height
                img_bbox = image_data["bbox"]
                
                # Scale image to fit slide properly
                max_width = prs.slide_width - Inches(1.0)
                max_height = prs.slide_height - Inches(2.0)
                
                # Calculate aspect ratio
                img_width = img_bbox[2] - img_bbox[0]
                img_height = img_bbox[3] - img_bbox[1]
                
                if img_width > 0 and img_height > 0:
                    aspect_ratio = img_height / img_width
                    
                    # Fit image to slide
                    if max_width * aspect_ratio <= max_height:
                        display_width = max_width
                        display_height = max_width * aspect_ratio
                    else:
                        display_height = max_height
                        display_width = max_height / aspect_ratio
                    
                    # Center the image
                    left = (prs.slide_width - display_width) / 2
                    top = Inches(1.5) + (max_height - display_height) / 2
                    
                    # Add image to slide
                    image_stream = io.BytesIO(image_data["bytes"])
                    slide.shapes.add_picture(image_stream, left, top, display_width, display_height)
                
            except Exception as e:
                # Add error message if image fails
                error_box = slide.shapes.add_textbox(
                    Inches(2), Inches(3), Inches(8), Inches(2)
                )
                error_frame = error_box.text_frame
                p = error_frame.paragraphs[0]
                p.text = f"Image could not be displayed\nError: {str(e)[:50]}..."
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.CENTER

def create_pdf_content_slides(prs, pdf_pages, section_title):
    """
    Wrapper function that calls the improved PDF slide creation.
    """
    create_pdf_content_slides_improved(prs, pdf_pages, section_title)

def create_excel_slides(prs, excel_df, section_title):
    """
    Creates slides for Excel data with pagination.
    """
    if excel_df.empty:
        return
    
    blank_layout = prs.slide_layouts[6]
    max_rows_per_slide = 20
    total_rows = len(excel_df)
    num_slides = (total_rows + max_rows_per_slide - 1) // max_rows_per_slide

    for slide_idx in range(num_slides):
        slide = prs.slides.add_slide(blank_layout)
        title_text = section_title
        if num_slides > 1:
            title_text += f" (Page {slide_idx + 1}/{num_slides})"
        add_slide_banner(slide, title_text)

        start_row = slide_idx * max_rows_per_slide
        end_row = min(start_row + max_rows_per_slide, total_rows)
        df_slice = excel_df.iloc[start_row:end_row]

        rows, cols = df_slice.shape
        rows += 1

        table_left = Inches(0.5)
        table_top = Inches(1.2)
        table_width = prs.slide_width - Inches(1.0)
        table_height = prs.slide_height - Inches(1.5)

        table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
        table = table_shape.table

        for c in range(cols):
            table.columns[c].width = int(table_width / cols)

        for c, col_name in enumerate(df_slice.columns):
            cell = table.cell(0, c)
            p = cell.text_frame.paragraphs[0]
            p.text = str(col_name)
            p.font.bold = True
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(138, 43, 226)
            p.font.color.rgb = RGBColor(255, 255, 255)

        for r, row_data in enumerate(df_slice.itertuples(index=False)):
            for c, value in enumerate(row_data):
                cell = table.cell(r + 1, c)
                p = cell.text_frame.paragraphs[0]
                p.text = str(value) if pd.notna(value) else ""
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER


def create_cad_slides(prs, cad_image_paths, section_title):
    """
    Creates a slide for each CAD screenshot, ensuring it uses the full content width.
    """
    slide_layout = prs.slide_layouts[6]
    for i, cad_path in enumerate(cad_image_paths):
        if cad_path and os.path.exists(cad_path):
            slide = prs.slides.add_slide(slide_layout)
            add_slide_banner(slide, f"{section_title} - Drawing {i+1}")
            
            content_area_top = Inches(0.8)
            content_area_height = prs.slide_height - content_area_top
            
            try:
                display_w = prs.slide_width
                
                img = PILImage.open(cad_path)
                aspect = img.height / img.width
                display_h = display_w * aspect

                left = 0
                top = content_area_top + (content_area_height - display_h) / 2
                
                if top < content_area_top:
                    top = content_area_top
                    display_h = content_area_height
                    display_w = display_h / aspect
                    left = (prs.slide_width - display_w) / 2

                slide.shapes.add_picture(cad_path, left, top, width=display_w)
            except Exception as e:
                st.error(f"Error adding CAD image {i+1}: {e}")

def create_nx_model_slides(prs, nx_model_groups, section_title):
    """
    Creates slides for ALL NX models, one slide per model.
    """
    slide_layout = prs.slide_layouts[6]
    
    # Create a slide for each model group
    for group_idx, model_group in enumerate(nx_model_groups):
        if len(model_group) < 3: 
            continue

        slide = prs.slides.add_slide(slide_layout)
        add_slide_banner(slide, f"{section_title} - Model {group_idx + 1}")
        
        total_content_width = prs.slide_width - Inches(1.0)
        gap = Inches(0.15)
        
        image_width = (total_content_width - (2 * gap)) / 3
        
        view_paths = model_group[:3]
        view_names = ['Front View', 'Back View', 'Isometric View']
        
        current_left = Inches(0.5)
        
        for i, img_path in enumerate(view_paths):
            try:
                original_img = PILImage.open(img_path)
                centered_img = center_tire_properly(original_img)
                
                # Create unique centered image path
                centered_path = img_path.replace('.png', f'_centered_slide{group_idx}.png')
                centered_img.save(centered_path)
                
                aspect_ratio = centered_img.height / centered_img.width
                image_height = image_width * aspect_ratio

                top = Inches(1.0) + ((prs.slide_height - Inches(1.0) - image_height) / 2)

                slide.shapes.add_picture(centered_path, current_left, top, width=image_width, height=image_height)
                
                label_top = top + image_height + Inches(0.1)
                label_box = slide.shapes.add_textbox(current_left, label_top, image_width, Inches(0.3))
                label_frame = label_box.text_frame
                p = label_frame.paragraphs[0]
                p.text = view_names[i] if i < len(view_names) else f"View {i+1}"
                p.font.size = Pt(14)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
                current_left += image_width + gap

            except Exception as e:
                st.error(f"Could not add NX image {i+1} to slide: {e}")

def create_nx_closeup_slide(prs, nx_model_groups, section_title):
    """
    Creates close-up slides for ALL models (one slide per model).
    """
    if not nx_model_groups or len(nx_model_groups) == 0:
        return
    
    slide_layout = prs.slide_layouts[6]
    
    # Create close-up slide for each model
    for group_idx, model_group in enumerate(nx_model_groups):
        if len(model_group) < 3:
            continue
            
        slide = prs.slides.add_slide(slide_layout)
        
        # Update title to include model number
        if len(nx_model_groups) > 1:
            slide_title = f"{section_title} - Model {group_idx + 1}"
        else:
            slide_title = section_title
            
        add_slide_banner(slide, slide_title)
        
        front_view_path = model_group[1]  # Front View
        isometric_view_path = model_group[2]  # Isometric View
        
        gap = Inches(0.5)
        image_width = (prs.slide_width - Inches(1.0) - gap) / 2
        
        left_img_left = Inches(0.5)
        right_img_left = left_img_left + image_width + gap
        
        content_top = Inches(1.2)
        available_height = prs.slide_height - content_top - Inches(1.0)
        
        def create_zoomed_centered_view(img_path, zoom_factor=2.5, model_idx=0):
            """Create zoomed centered view with unique filename"""
            try:
                img = PILImage.open(img_path)
                centered_img = center_tire_properly(img)
                
                width, height = centered_img.size
                
                crop_width = int(width / zoom_factor)
                crop_height = int(height / zoom_factor)
                
                center_x = width // 2
                center_y = height // 2
                
                left = center_x - crop_width // 2
                top = center_y - crop_height // 2
                right = left + crop_width
                bottom = top + crop_height
                
                left = max(0, left)
                top = max(0, top)
                right = min(width, right)
                bottom = min(height, bottom)
                
                cropped = centered_img.crop((left, top, right, bottom))
                resized = cropped.resize((width, height), PILImage.Resampling.LANCZOS)
                
                # Create unique filename for each model
                processed_path = img_path.replace('.png', f'_zoomed_model{model_idx}.png')
                resized.save(processed_path)
                return processed_path
            except Exception:
                return img_path
        
        # Front view - with model-specific filename
        try:
            if os.path.exists(front_view_path):
                zoomed_front_path = create_zoomed_centered_view(front_view_path, zoom_factor=2.5, model_idx=group_idx)
                
                img = PILImage.open(zoomed_front_path)
                aspect_ratio = img.height / img.width
                image_height = min(image_width * aspect_ratio, available_height - Inches(0.5))
                
                top_pos = content_top + (available_height - image_height - Inches(0.5)) / 2
                
                slide.shapes.add_picture(zoomed_front_path, left_img_left, top_pos, 
                                       width=image_width, height=image_height)
                
                label_top = top_pos + image_height + Inches(0.1)
                label_box = slide.shapes.add_textbox(left_img_left, label_top, image_width, Inches(0.4))
                label_frame = label_box.text_frame
                p = label_frame.paragraphs[0]
                p.text = "Front View"
                p.font.size = Pt(16)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
        except Exception as e:
            st.error(f"Error adding front view for model {group_idx + 1}: {e}")
        
        # Isometric view - with model-specific filename
        try:
            if os.path.exists(isometric_view_path):
                zoomed_iso_path = create_zoomed_centered_view(isometric_view_path, zoom_factor=2.5, model_idx=group_idx)
                
                img = PILImage.open(zoomed_iso_path)
                aspect_ratio = img.height / img.width
                image_height = min(image_width * aspect_ratio, available_height - Inches(0.5))
                
                top_pos = content_top + (available_height - image_height - Inches(0.5)) / 2
                
                slide.shapes.add_picture(zoomed_iso_path, right_img_left, top_pos, 
                                       width=image_width, height=image_height)
                
                label_top = top_pos + image_height + Inches(0.1)
                label_box = slide.shapes.add_textbox(right_img_left, label_top, image_width, Inches(0.4))
                label_frame = label_box.text_frame
                p = label_frame.paragraphs[0]
                p.text = "Isometric View"
                p.font.size = Pt(16)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
        except Exception as e:
            st.error(f"Error adding isometric view for model {group_idx + 1}: {e}")


# ---------- STREAMLIT GUI ----------
def main():
    """
    The main function that runs the Streamlit application.
    """
    add_logo_to_streamlit()
    st.title("Enhanced Tyre Report Generator")
    initialize_session_state()

    tab1, tab2, tab3 = st.tabs([
        "📸 Capture CAD Drawing", 
        "📐 Capture NX 3D Models", 
        "📄 Generate Reports"
    ])

    with tab1:
        updated_tab1_section()
    with tab2:
        updated_tab2_section()
    with tab3:
        updated_tab3_section()

def updated_tab1_section():
    """
    UI and logic for the 'Capture CAD Drawing' tab.
    """
    st.header("📸 AutoCAD Drawing Capture")
    cad_files = st.file_uploader("Upload 2D CAD Files (.dwg)", type=["dwg"], accept_multiple_files=True)
    
    if cad_files:
        if st.button("🚀 Process All CAD Files"):
            paths = []
            with st.spinner('Processing CAD files... This is now faster.'):
                for i, cad_file in enumerate(cad_files):
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".dwg") as tmp:
                        tmp.write(cad_file.getvalue())
                        cad_path = tmp.name
                    
                    ss_path = open_autocad_and_capture_screenshot(cad_path)
                    
                    try:
                        os.remove(cad_path)
                    except OSError as e:
                        st.warning(f"Could not remove temp file {cad_path}: {e}")

                    if ss_path:
                        final_path = os.path.join(tempfile.gettempdir(), f"cad_ss_{i}.png")
                        shutil.copy(ss_path, final_path)
                        paths.append(final_path)
                        st.success(f"✅ Captured {cad_file.name}")
                    else:
                        st.error(f"❌ Failed to capture {cad_file.name}")

            if paths:
                st.session_state['cad_screenshot_paths'] = paths
                st.session_state['cad_screenshots_captured'] = True
                st.success("🎉 All CAD drawings captured!")
                for path in paths:
                    st.image(path)

def updated_tab2_section():
    """
    Simple NX UI with single button to capture multiple models.
    """
    st.header("📐 NX 3D Model Capture")
    
    if st.session_state.get('nx_model_groups'):
        st.success(f"✅ {len(st.session_state.get('nx_model_groups', []))} NX model(s) captured.")
        with st.expander("View Captured Models"):
            for i, model_group in enumerate(st.session_state.get('nx_model_groups', [])):
                st.subheader(f"Model {i+1}")
                cols = st.columns(len(model_group))
                for col, path in zip(cols, model_group):
                    col.image(path)

    if st.session_state.get('nx_screenshots_captured'):
        if st.button("🗑️ Clear All NX Captures"):
            st.session_state['nx_model_groups'] = []
            st.session_state['nx_screenshots_captured'] = False
            st.rerun()

    if st.button("🚀 Capture New NX Model"):
        with st.spinner("Waiting for NX automation..."):
            views = open_nx_and_capture_views_manual()
            if views and len(views) == 3:
                model_group = [views[v] for v in ['Top View', 'Front View', 'Isometric View']]
                st.session_state.setdefault('nx_model_groups', []).append(model_group)
                st.session_state['nx_screenshots_captured'] = True
                st.success("✅ Capture complete! Please return to the Streamlit app.")
                st.rerun()
            else:
                st.error("❌ Failed to capture all 3 required views from NX.")

def updated_tab3_section():
    """
    UI and logic for the 'Generate Reports' tab.
    """
    st.header("📄 Generate Enhanced Reports")
    
    with st.expander("🔧 Software Configuration (Optional)"):
        st.info("The app tries to auto-detect software. Use these fields to override.")
        custom_autocad = st.text_input("Custom AutoCAD Path:", st.session_state.get('custom_autocad_path', ''))
        st.session_state['custom_autocad_path'] = custom_autocad
        
        custom_nx = st.text_input("Custom NX Path:", st.session_state.get('custom_nx_path', ''))
        st.session_state['custom_nx_path'] = custom_nx

    pdf_files = st.file_uploader("Upload PDF Files", type=["pdf"], accept_multiple_files=True)
    excel_files = st.file_uploader("Upload Excel Files", type=["xlsx"], accept_multiple_files=True)
    
    st.info(f"CAD Drawings: {'✅ Ready' if st.session_state.get('cad_screenshots_captured') else '❌ Not captured'}")
    st.info(f"NX Models: {'✅ Ready' if st.session_state.get('nx_screenshots_captured') else '❌ Not captured'}")
        
    if st.button("🚀 Generate Enhanced PowerPoint Report"):
        if not all([pdf_files, excel_files, st.session_state.get('cad_screenshots_captured'), st.session_state.get('nx_screenshots_captured')]):
            st.error("Missing one or more inputs. Please provide all files and captures.")
        else:
            with st.spinner("Generating report... This may take a moment."):
                with tempfile.TemporaryDirectory() as tmpdir:
                    all_pdf_pages = []
                    for pdf_file in pdf_files:
                        pdf_path = os.path.join(tmpdir, pdf_file.name)
                        with open(pdf_path, "wb") as f:
                            f.write(pdf_file.getvalue())
                        all_pdf_pages.extend(extract_pdf_elements_improved(pdf_path))
                    
                    all_excel_df = pd.concat(
                        [read_excel_data(excel_file) for excel_file in excel_files], 
                        ignore_index=True
                    )
                    
                    ppt_out_path = os.path.join(tmpdir, "Enhanced_Tyre_Report.pptx")
                    
                    generate_ppt(
                        all_pdf_pages, 
                        all_excel_df, 
                        st.session_state.get('cad_screenshot_paths', []), 
                        st.session_state.get('nx_model_groups', []), 
                        ppt_out_path
                    )
                    
                    with open(ppt_out_path, "rb") as f:
                        st.download_button(
                            "📊 Download Enhanced PowerPoint Report", 
                            f, 
                            file_name="Enhanced_Tyre_Report.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
            st.success("Report generated!")

if __name__ == "__main__":
    main()